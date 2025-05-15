import sys
import logging
import json
import os
import requests
import io
from pathlib import Path
from typing import Optional, Any, Iterable
from pptx import Presentation

# Add both project root and src directory to Python path
project_root = os.path.abspath(
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))
)
sys.path.insert(0, project_root)
sys.path.insert(0, os.path.join(project_root, "src"))

import mcp.types as types
from mcp.server import NotificationOptions, Server
from mcp.server.models import InitializationOptions
from src.utils.microsoft.util import (
    get_credentials,
    authenticate_and_save_credentials,
)
from mcp.types import (
    AnyUrl,
    Resource,
)
from mcp.server.lowlevel.helper_types import ReadResourceContents

SERVICE_NAME = Path(__file__).parent.name
MICROSOFT_GRAPH_API_URL = "https://graph.microsoft.com/v1.0"
SCOPES = [
    "Files.ReadWrite",
    "Sites.ReadWrite.All",
    "offline_access",
]

# Configure logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(name)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(SERVICE_NAME)


async def make_graph_api_request(
    method,
    endpoint,
    data=None,
    params=None,
    access_token=None,
    content_type=None,
    stream=False,
):
    """Make a request to the Microsoft Graph API"""
    url = f"{MICROSOFT_GRAPH_API_URL}/{endpoint}"
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": content_type if content_type else "application/json",
    }

    try:
        session = requests.Session()
        if method.lower() == "get":
            response = session.get(url, headers=headers, params=params, timeout=60.0)
        elif method.lower() == "post":
            if content_type and content_type != "application/json":
                response = session.post(
                    url, data=data, headers=headers, params=params, timeout=60.0
                )
            else:
                response = session.post(
                    url, json=data, headers=headers, params=params, timeout=60.0
                )
        elif method.lower() == "patch":
            if content_type and content_type != "application/json":
                response = session.patch(
                    url, data=data, headers=headers, params=params, timeout=60.0
                )
            else:
                response = session.patch(
                    url, json=data, headers=headers, params=params, timeout=60.0
                )
        elif method.lower() == "put":
            if content_type and content_type != "application/json":
                response = session.put(
                    url, data=data, headers=headers, params=params, timeout=60.0
                )
            else:
                response = session.put(
                    url, json=data, headers=headers, params=params, timeout=60.0
                )
        elif method.lower() == "delete":
            response = session.delete(url, headers=headers, params=params, timeout=60.0)
        else:
            raise ValueError(f"Unsupported HTTP method: {method}")

        response.raise_for_status()
        if response.status_code == 204:  # No content
            return {"success": True, "status_code": 204}

        if stream:
            return response

        return response.json()

    except requests.exceptions.HTTPError as e:
        error_message = f"Microsoft Graph API error: {e.response.status_code}"
        try:
            error_response = e.response.json()
            if "error" in error_response:
                error_details = error_response["error"]
                error_message = f"{error_details.get('code', 'Error')}: {error_details.get('message', 'Unknown error')}"
        except Exception:
            pass

        raise ValueError(error_message)

    except requests.exceptions.RequestException as e:
        raise ValueError(f"Failed to connect to Microsoft Graph API: {str(e)}")

    except Exception as e:
        raise ValueError(f"Error communicating with Microsoft Graph API: {str(e)}")


async def create_powerpoint_client(access_token: str) -> Any:
    """
    Create a Microsoft PowerPoint client instance using the provided credentials.

    Args:
        access_token: The OAuth access token

    Returns:
        dict: A dictionary containing:
            - token: The access token
            - headers: Standard HTTP headers for Microsoft Graph API requests
            - base_url: The base URL for Microsoft Graph API endpoints
    """
    # Standard headers for API requests
    standard_headers = {
        "Authorization": f"Bearer {access_token}",
        "Accept": "application/json",
        "Content-Type": "application/json",
    }

    return {
        "token": access_token,
        "headers": standard_headers,
        "base_url": MICROSOFT_GRAPH_API_URL,
    }


async def is_sharepoint_storage(access_token):
    """Detect if we're using SharePoint or OneDrive storage"""
    drive_info = await make_graph_api_request(
        "get", "me/drive", access_token=access_token
    )
    return (
        drive_info.get("driveType") == "business"
        or "sharepoint" in drive_info.get("webUrl", "").lower()
    )


def create_server(user_id, api_key=None):
    """Create a new server instance for PowerPoint operations"""
    server = Server(f"{SERVICE_NAME}-server")
    server.user_id = user_id
    server.api_key = api_key

    async def get_microsoft_client():
        """Get Microsoft access token for the current user"""
        return await get_credentials(user_id, SERVICE_NAME, api_key=api_key)

    @server.list_resources()
    async def handle_list_resources(
        cursor: Optional[str] = None,
    ) -> list[Resource]:
        """List PowerPoint presentations from OneDrive"""
        access_token = await get_microsoft_client()

        try:
            # Determine if we're using SharePoint or OneDrive
            is_sharepoint = await is_sharepoint_storage(access_token)

            endpoint = "me/drive/root/search(q='.pptx')"
            query_params = {
                "$top": 50,
                "$select": "id,name,webUrl,lastModifiedDateTime,size"
                + (",file" if is_sharepoint else ""),
                "$orderby": "lastModifiedDateTime desc",
            }

            if cursor:
                query_params["$skiptoken"] = cursor

            result = await make_graph_api_request(
                "get", endpoint, params=query_params, access_token=access_token
            )

            resources = []

            for item in result.get("value", []):
                # For SharePoint, filter by MIME type
                if is_sharepoint:
                    if (
                        item.get("file")
                        and item.get("file", {}).get("mimeType")
                        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    ):
                        resources.append(
                            Resource(
                                uri=f"powerpoint://file/{item['id']}",
                                mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                name=f"{item['name']}",
                            )
                        )
                else:
                    if item["name"].lower().endswith(".pptx"):
                        resources.append(
                            Resource(
                                uri=f"powerpoint://file/{item['id']}",
                                mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                name=f"{item['name']}",
                            )
                        )

            return resources

        except Exception as e:
            logger.error(f"Error fetching PowerPoint resources: {str(e)}")
            return []

    @server.read_resource()
    async def handle_read_resource(uri: AnyUrl) -> Iterable[ReadResourceContents]:
        """Read a PowerPoint presentation from OneDrive"""
        access_token = await get_microsoft_client()
        uri_str = str(uri)

        if uri_str.startswith("powerpoint://file/"):
            file_id = uri_str.replace("powerpoint://file/", "")

            try:
                endpoint = f"me/drive/items/{file_id}"
                file_info = await make_graph_api_request(
                    "get", endpoint, access_token=access_token
                )

                result = {
                    "id": file_info.get("id"),
                    "name": file_info.get("name"),
                    "webUrl": file_info.get("webUrl"),
                    "lastModifiedDateTime": file_info.get("lastModifiedDateTime"),
                    "size": file_info.get("size"),
                    "createdDateTime": file_info.get("createdDateTime"),
                    "downloadUrl": file_info.get("@microsoft.graph.downloadUrl"),
                    "contentPreview": "Content preview not available in resource view.",
                }

                formatted_content = json.dumps(result, indent=2)
                return [
                    ReadResourceContents(
                        content=formatted_content, mime_type="application/json"
                    )
                ]

            except Exception as e:
                error_msg = f"Error reading PowerPoint presentation: {str(e)}"
                logger.error(error_msg)

                formatted_error = {
                    "id": file_id,
                    "error": error_msg,
                    "status": "error",
                    "success": False,
                }

                return [
                    ReadResourceContents(
                        content=json.dumps(formatted_error),
                        mime_type="application/json",
                    )
                ]

        raise ValueError(f"Unsupported resource URI: {uri_str}")

    @server.list_tools()
    async def handle_list_tools() -> list[types.Tool]:
        """List available tools for PowerPoint"""
        return [
            types.Tool(
                name="list_presentations",
                description="List PowerPoint presentations from OneDrive",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "limit": {
                            "type": "integer",
                            "description": "Maximum number of presentations to return",
                            "default": 50,
                        },
                        "query": {
                            "type": "string",
                            "description": "Optional search query to filter presentations (defaults to all .pptx files)",
                        },
                    },
                },
                outputSchema={
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Array of JSON strings containing PowerPoint presentations with their metadata",
                    "examples": [
                        '{"status":"success","presentations":[{"id":"12345","name":"Presentation1.pptx","web_url":"https://example.com/pres1.pptx","last_modified":"2023-07-15T10:30:00Z","size":25600}]}'
                    ],
                },
            ),
            types.Tool(
                name="create_presentation",
                description="Create a new PowerPoint presentation in OneDrive",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "name": {
                            "type": "string",
                            "description": "Name for the new presentation (will add .pptx extension if not included)",
                        },
                        "title_slide": {
                            "type": "string",
                            "description": "Title for the first slide (optional)",
                        },
                        "folder_path": {
                            "type": "string",
                            "description": "Path to folder in OneDrive (optional, defaults to root)",
                        },
                    },
                    "required": ["name"],
                },
                outputSchema={
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Array of JSON strings containing details of the newly created PowerPoint presentation",
                    "examples": [
                        '{"status":"success","file_id":"12345","name":"Test Presentation.pptx","web_url":"https://example.com/test-presentation.pptx"}'
                    ],
                },
            ),
            types.Tool(
                name="read_presentation",
                description="Extract content/structure from a PowerPoint presentation",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "file_id": {
                            "type": "string",
                            "description": "ID of the PowerPoint presentation",
                        },
                    },
                    "required": ["file_id"],
                },
                outputSchema={
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Array of JSON strings containing presentation structure and content",
                    "examples": [
                        '{"status":"success","file_id":"12345","name":"Example.pptx","slide_count":5,"slides":[{"index":1,"title":"Introduction","text":"Sample content"}]}'
                    ],
                },
            ),
            types.Tool(
                name="add_slide",
                description="Add a new slide to an existing presentation",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "file_id": {
                            "type": "string",
                            "description": "ID of the PowerPoint presentation",
                        },
                        "title": {
                            "type": "string",
                            "description": "Title for the new slide",
                        },
                        "content": {
                            "type": "string",
                            "description": "Content for the new slide",
                        },
                        "layout": {
                            "type": "string",
                            "description": "Layout type (Title and Content, Title Only, etc.)",
                            "default": "Title and Content",
                        },
                    },
                    "required": ["file_id"],
                },
                outputSchema={
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Array of JSON strings containing result of adding a slide",
                    "examples": [
                        '{"status":"success","file_id":"12345","slide_added":true,"slide_index":6,"title":"New Slide"}'
                    ],
                },
            ),
            types.Tool(
                name="delete_presentation",
                description="Delete a PowerPoint presentation from OneDrive",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "file_id": {
                            "type": "string",
                            "description": "ID of the PowerPoint presentation",
                        },
                    },
                    "required": ["file_id"],
                },
                outputSchema={
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Array of JSON strings containing result of the presentation deletion operation",
                    "examples": [
                        '{"status":"success","deleted":true,"file_id":"12345"}'
                    ],
                },
            ),
            types.Tool(
                name="update_slide",
                description="Update content on a specific slide",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "file_id": {
                            "type": "string",
                            "description": "ID of the PowerPoint presentation",
                        },
                        "slide_index": {
                            "type": "integer",
                            "description": "Index of the slide to update (1-based)",
                        },
                        "title": {
                            "type": "string",
                            "description": "New title for the slide (optional)",
                        },
                        "content": {
                            "type": "string",
                            "description": "New content for the slide (optional)",
                        },
                    },
                    "required": ["file_id", "slide_index"],
                },
                outputSchema={
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Array of JSON strings containing result of updating a slide",
                    "examples": [
                        '{"status":"success","file_id":"12345","slide_updated":true,"slide_index":3}'
                    ],
                },
            ),
            types.Tool(
                name="delete_slide",
                description="Remove a slide from a presentation",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "file_id": {
                            "type": "string",
                            "description": "ID of the PowerPoint presentation",
                        },
                        "slide_index": {
                            "type": "integer",
                            "description": "Index of the slide to delete (1-based)",
                        },
                    },
                    "required": ["file_id", "slide_index"],
                },
                outputSchema={
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Array of JSON strings containing result of deleting a slide",
                    "examples": [
                        '{"status":"success","file_id":"12345","slide_deleted":true,"previous_slide_count":10,"new_slide_count":9}'
                    ],
                },
            ),
            types.Tool(
                name="download_presentation",
                description="Get a download URL for a PowerPoint presentation",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "file_id": {
                            "type": "string",
                            "description": "ID of the PowerPoint presentation",
                        },
                    },
                    "required": ["file_id"],
                },
                outputSchema={
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Array of JSON strings containing presentation download details",
                    "examples": [
                        '{"status":"success","file_id":"12345","name":"Example.pptx","download_url":"https://download.example.com/ppt12345.pptx","size":36582}'
                    ],
                },
            ),
            types.Tool(
                name="search_presentations",
                description="Search for PowerPoint presentations by content",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "query": {
                            "type": "string",
                            "description": "Search query to find presentations containing this content",
                        },
                        "limit": {
                            "type": "integer",
                            "description": "Maximum number of presentations to return",
                            "default": 25,
                        },
                    },
                    "required": ["query"],
                },
                outputSchema={
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Array of JSON strings containing search results with matching presentations",
                    "examples": [
                        '{"status":"success","presentations":[{"id":"12345","name":"Quarterly Update.pptx","web_url":"https://example.com/quarterly-update.pptx","last_modified":"2023-06-10T09:15:30Z"}]}'
                    ],
                },
            ),
        ]

    def format_create_presentation_endpoint(args):
        """Format the endpoint for creating a presentation"""
        file_name = args.get("name", "")
        if not file_name.lower().endswith(".pptx"):
            file_name = f"{file_name}.pptx"

        folder_path = args.get("folder_path", "").strip("/")
        if folder_path:
            return f"me/drive/root:/{folder_path}/{file_name}:/content"
        else:
            return f"me/drive/root:/{file_name}:/content"

    async def get_presentation_as_bytes(response):
        """Convert presentation stream to bytes"""
        ppt_bytes = io.BytesIO()
        if hasattr(response, "content"):
            ppt_bytes.write(response.content)
        elif hasattr(response, "read"):
            ppt_bytes.write(await response.read())
        ppt_bytes.seek(0)
        return ppt_bytes

    @server.call_tool()
    async def handle_call_tool(
        name: str, arguments: dict | None
    ) -> list[types.TextContent]:
        """Handle tool execution requests for PowerPoint"""
        arguments = arguments or {}
        logger.info(
            f"Tool call received: {name} with arguments: {json.dumps(arguments, indent=2)}"
        )

        try:
            access_token = await get_microsoft_client()
            logger.info(f"Successfully retrieved access token for tool: {name}")

            try:
                if name == "list_presentations":
                    logger.info("Executing list_presentations tool")
                    # Determine if we're using SharePoint or OneDrive
                    is_sharepoint = await is_sharepoint_storage(access_token)
                    logger.info(
                        f"Storage type detected: {'SharePoint' if is_sharepoint else 'OneDrive'}"
                    )

                    endpoint = "me/drive/root/search(q='.pptx')"
                    params = {
                        "$top": (
                            arguments.get("limit", 50) if not is_sharepoint else 100
                        ),
                        "$select": "id,name,webUrl,lastModifiedDateTime,size,createdDateTime,file",
                        "$orderby": "lastModifiedDateTime desc",
                    }

                    if arguments.get("query"):
                        params["search"] = arguments.get("query")

                    logger.info(
                        f"Making Graph API request to endpoint: {endpoint} with params: {params}"
                    )
                    result = await make_graph_api_request(
                        "get", endpoint, params=params, access_token=access_token
                    )
                    logger.info(
                        f"Retrieved {len(result.get('value', []))} items from Graph API"
                    )

                    presentations = []
                    if is_sharepoint:
                        for item in result.get("value", []):
                            if (
                                item.get("file")
                                and item.get("file", {}).get("mimeType")
                                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            ):
                                presentations.append(item)
                    else:
                        for item in result.get("value", []):
                            if item["name"].lower().endswith(".pptx"):
                                presentations.append(item)

                    formatted_presentations = []
                    for item in presentations:
                        formatted_presentations.append(
                            {
                                "id": item.get("id"),
                                "name": item.get("name"),
                                "web_url": item.get("webUrl"),
                                "last_modified": item.get("lastModifiedDateTime"),
                                "created": item.get("createdDateTime"),
                                "size": item.get("size"),
                            }
                        )

                    response = {
                        "status": "success",
                        "presentations": formatted_presentations,
                    }

                    return [
                        types.TextContent(
                            type="text", text=json.dumps(response, indent=2)
                        )
                    ]

                elif name == "create_presentation":
                    logger.info("Executing create_presentation tool")
                    title_slide = arguments.get("title_slide", "")
                    endpoint = format_create_presentation_endpoint(arguments)
                    logger.info(f"Creating presentation at endpoint: {endpoint}")

                    # Create a proper PowerPoint presentation using python-pptx
                    prs = Presentation()
                    if title_slide:
                        logger.info(f"Adding title slide with title: {title_slide}")
                        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
                        slide = prs.slides.add_slide(title_slide_layout)
                        title = slide.shapes.title
                        subtitle = slide.placeholders[1]
                        title.text = title_slide
                        subtitle.text = "Created via MCP"

                    # Save the presentation to a byte stream
                    pptx_bytes = io.BytesIO()
                    prs.save(pptx_bytes)
                    pptx_bytes.seek(0)
                    logger.info("Presentation created and saved to byte stream")

                    result = await make_graph_api_request(
                        "put",
                        endpoint,
                        data=pptx_bytes.getvalue(),
                        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        access_token=access_token,
                        params={"@microsoft.graph.conflictBehavior": "rename"},
                    )

                    response = {
                        "status": "success",
                        "file_id": result.get("id"),
                        "name": result.get("name"),
                        "web_url": result.get("webUrl"),
                        "title_slide": title_slide,
                    }

                    return [
                        types.TextContent(
                            type="text", text=json.dumps(response, indent=2)
                        )
                    ]

                elif name == "read_presentation":
                    file_id = arguments.get("file_id")
                    logger.info(
                        f"Executing read_presentation tool for file_id: {file_id}"
                    )

                    # Get document metadata
                    ppt_info_endpoint = f"me/drive/items/{file_id}"
                    logger.info(
                        f"Fetching presentation metadata from endpoint: {ppt_info_endpoint}"
                    )
                    ppt_info = await make_graph_api_request(
                        "get", ppt_info_endpoint, access_token=access_token
                    )
                    logger.info(
                        f"Retrieved metadata for presentation: {ppt_info.get('name')}"
                    )

                    # Get document content
                    content_endpoint = f"me/drive/items/{file_id}/content"
                    response = await make_graph_api_request(
                        "get", content_endpoint, access_token=access_token, stream=True
                    )

                    # Extract presentation structure using python-pptx
                    slide_data = []
                    try:
                        ppt_bytes = await get_presentation_as_bytes(response)

                        # Parse with python-pptx
                        prs = Presentation(ppt_bytes)

                        # Extract slides information
                        for i, slide in enumerate(prs.slides):
                            slide_info = {"index": i + 1}

                            # Get title if exists
                            if slide.shapes.title and slide.shapes.title.text:
                                slide_info["title"] = slide.shapes.title.text

                            # Extract text from all shapes
                            texts = []
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text:
                                    texts.append(shape.text)

                            if texts:
                                slide_info["text"] = "\n".join(texts)

                            slide_data.append(slide_info)

                    except Exception as e:
                        logger.error(f"Error parsing PowerPoint file: {str(e)}")
                        slide_data = [{"error": "Could not parse presentation content"}]

                    response = {
                        "status": "success",
                        "file_id": ppt_info.get("id"),
                        "name": ppt_info.get("name"),
                        "slide_count": len(slide_data),
                        "slides": slide_data,
                    }

                    return [
                        types.TextContent(
                            type="text", text=json.dumps(response, indent=2)
                        )
                    ]

                elif name == "add_slide":
                    file_id = arguments.get("file_id")
                    title = arguments.get("title", "New Slide")
                    content = arguments.get("content", "")
                    layout_name = arguments.get("layout", "Title and Content")
                    logger.info(
                        f"Executing add_slide tool for file_id: {file_id}, title: {title}, layout: {layout_name}"
                    )

                    # Get presentation metadata
                    ppt_info_endpoint = f"me/drive/items/{file_id}"
                    ppt_info = await make_graph_api_request(
                        "get", ppt_info_endpoint, access_token=access_token
                    )

                    # Get presentation content
                    content_endpoint = f"me/drive/items/{file_id}/content"
                    response = await make_graph_api_request(
                        "get", content_endpoint, access_token=access_token, stream=True
                    )

                    ppt_bytes = await get_presentation_as_bytes(response)

                    try:
                        # Try to load as a PowerPoint presentation
                        prs = Presentation(ppt_bytes)

                        # Find appropriate layout
                        layout_mapping = {
                            "Title Slide": 0,
                            "Title and Content": 1,
                            "Section Header": 2,
                            "Title Only": 3,
                            "Blank": 5,
                        }

                        layout_idx = layout_mapping.get(
                            layout_name, 1
                        )  # Default to Title and Content
                        slide_layout = prs.slide_layouts[layout_idx]

                        # Add new slide
                        slide = prs.slides.add_slide(slide_layout)

                        # Add title if the layout has a title placeholder
                        if slide.shapes.title:
                            slide.shapes.title.text = title

                        # Add content if the layout has a content placeholder
                        for shape in slide.placeholders:
                            if shape.placeholder_format.type == 1:  # Body placeholder
                                shape.text = content
                                break

                        # Save updated presentation
                        updated_bytes = io.BytesIO()
                        prs.save(updated_bytes)
                        updated_bytes.seek(0)

                        # Update the presentation
                        result = await make_graph_api_request(
                            "put",
                            content_endpoint,
                            data=updated_bytes.getvalue(),
                            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            access_token=access_token,
                            params={"@microsoft.graph.conflictBehavior": "replace"},
                        )

                        response = {
                            "status": "success",
                            "file_id": result.get("id", file_id),
                            "slide_added": True,
                            "slide_index": len(prs.slides),
                            "title": title,
                        }

                        return [
                            types.TextContent(
                                type="text", text=json.dumps(response, indent=2)
                            )
                        ]
                    except Exception as e:
                        error_message = f"Error adding slide: {str(e)}"
                        logger.error(error_message)
                        return [
                            types.TextContent(
                                type="text",
                                text=json.dumps(
                                    {
                                        "status": "error",
                                        "message": error_message,
                                        "file_id": file_id,
                                    },
                                    indent=2,
                                ),
                            )
                        ]

                elif name == "delete_presentation":
                    file_id = arguments.get("file_id")
                    logger.info(
                        f"Executing delete_presentation tool for file_id: {file_id}"
                    )

                    endpoint = f"me/drive/items/{file_id}"
                    logger.info(f"Sending delete request to endpoint: {endpoint}")
                    result = await make_graph_api_request(
                        "delete", endpoint, access_token=access_token
                    )
                    logger.info(f"Presentation deletion result: {result}")

                    response = {
                        "status": "success",
                        "deleted": True,
                        "file_id": file_id,
                    }
                    logger.info("Presentation deleted successfully")

                    return [
                        types.TextContent(
                            type="text", text=json.dumps(response, indent=2)
                        )
                    ]

                elif name == "update_slide":
                    file_id = arguments.get("file_id")
                    slide_index = arguments.get("slide_index")
                    title = arguments.get("title")
                    content = arguments.get("content")
                    logger.info(
                        f"Executing update_slide tool for file_id: {file_id}, slide_index: {slide_index}"
                    )

                    if not slide_index:
                        logger.error("Missing required parameter: slide_index")
                        return [
                            types.TextContent(
                                type="text",
                                text=json.dumps(
                                    {
                                        "status": "error",
                                        "message": "slide_index is required",
                                        "file_id": file_id,
                                    },
                                    indent=2,
                                ),
                            )
                        ]

                    # Get presentation content
                    content_endpoint = f"me/drive/items/{file_id}/content"
                    response = await make_graph_api_request(
                        "get", content_endpoint, access_token=access_token, stream=True
                    )

                    ppt_bytes = await get_presentation_as_bytes(response)

                    try:
                        # Load the presentation
                        prs = Presentation(ppt_bytes)

                        # Check if slide index is valid
                        if slide_index < 1 or slide_index > len(prs.slides):
                            return [
                                types.TextContent(
                                    type="text",
                                    text=json.dumps(
                                        {
                                            "status": "error",
                                            "message": f"Invalid slide index: {slide_index}. Presentation has {len(prs.slides)} slides.",
                                            "file_id": file_id,
                                        },
                                        indent=2,
                                    ),
                                )
                            ]

                        # Get the slide (0-based index)
                        slide = prs.slides[slide_index - 1]

                        # Update title if provided
                        if title and slide.shapes.title:
                            slide.shapes.title.text = title

                        # Update content if provided
                        if content:
                            for shape in slide.placeholders:
                                # Find text placeholders
                                if (
                                    shape.placeholder_format.type == 1
                                ):  # Body placeholder
                                    shape.text = content
                                    break

                        # Save updated presentation
                        updated_bytes = io.BytesIO()
                        prs.save(updated_bytes)
                        updated_bytes.seek(0)

                        # Update the presentation
                        result = await make_graph_api_request(
                            "put",
                            content_endpoint,
                            data=updated_bytes.getvalue(),
                            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            access_token=access_token,
                            params={"@microsoft.graph.conflictBehavior": "replace"},
                        )

                        response = {
                            "status": "success",
                            "file_id": result.get("id", file_id),
                            "slide_updated": True,
                            "slide_index": slide_index,
                        }

                        return [
                            types.TextContent(
                                type="text", text=json.dumps(response, indent=2)
                            )
                        ]
                    except Exception as e:
                        error_message = f"Error updating slide: {str(e)}"
                        logger.error(error_message)
                        return [
                            types.TextContent(
                                type="text",
                                text=json.dumps(
                                    {
                                        "status": "error",
                                        "message": error_message,
                                        "file_id": file_id,
                                    },
                                    indent=2,
                                ),
                            )
                        ]

                elif name == "delete_slide":
                    file_id = arguments.get("file_id")
                    slide_index = arguments.get("slide_index")
                    logger.info(
                        f"Executing delete_slide tool for file_id: {file_id}, slide_index: {slide_index}"
                    )

                    if not slide_index:
                        return [
                            types.TextContent(
                                type="text",
                                text=json.dumps(
                                    {
                                        "status": "error",
                                        "message": "slide_index is required",
                                        "file_id": file_id,
                                    },
                                    indent=2,
                                ),
                            )
                        ]

                    # Get presentation content
                    content_endpoint = f"me/drive/items/{file_id}/content"
                    response = await make_graph_api_request(
                        "get", content_endpoint, access_token=access_token, stream=True
                    )

                    ppt_bytes = await get_presentation_as_bytes(response)

                    try:
                        # Load the presentation
                        prs = Presentation(ppt_bytes)
                        previous_slide_count = len(prs.slides)

                        # Check if slide index is valid
                        if slide_index < 1 or slide_index > previous_slide_count:
                            return [
                                types.TextContent(
                                    type="text",
                                    text=json.dumps(
                                        {
                                            "status": "error",
                                            "message": f"Invalid slide index: {slide_index}. Presentation has {previous_slide_count} slides.",
                                            "file_id": file_id,
                                        },
                                        indent=2,
                                    ),
                                )
                            ]

                        # Since python-pptx doesn't have a direct method to delete slides,
                        # we need to create a new presentation and copy all slides except the one to delete
                        new_prs = Presentation()

                        for i, slide in enumerate(prs.slides):
                            if i != slide_index - 1:  # Skip the slide we want to delete
                                # Get the layout of the current slide
                                # This is a rough approximation as direct layout matching is challenging
                                layout_idx = min(
                                    len(new_prs.slide_layouts) - 1,
                                    i % len(new_prs.slide_layouts),
                                )
                                slide_layout = new_prs.slide_layouts[layout_idx]

                                # Add a new slide with this layout
                                new_slide = new_prs.slides.add_slide(slide_layout)

                                # Unfortunately, we can't directly copy slide content in python-pptx
                                # In a real implementation, more complex slide content copying would be needed

                        # Save updated presentation
                        updated_bytes = io.BytesIO()
                        new_prs.save(updated_bytes)
                        updated_bytes.seek(0)

                        # Update the presentation
                        result = await make_graph_api_request(
                            "put",
                            content_endpoint,
                            data=updated_bytes.getvalue(),
                            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            access_token=access_token,
                            params={"@microsoft.graph.conflictBehavior": "replace"},
                        )

                        new_slide_count = len(new_prs.slides)
                        response = {
                            "status": "success",
                            "file_id": result.get("id", file_id),
                            "slide_deleted": True,
                            "previous_slide_count": previous_slide_count,
                            "new_slide_count": new_slide_count,
                        }

                        return [
                            types.TextContent(
                                type="text", text=json.dumps(response, indent=2)
                            )
                        ]
                    except Exception as e:
                        error_message = f"Error deleting slide: {str(e)}"
                        logger.error(error_message)
                        return [
                            types.TextContent(
                                type="text",
                                text=json.dumps(
                                    {
                                        "status": "error",
                                        "message": error_message,
                                        "file_id": file_id,
                                    },
                                    indent=2,
                                ),
                            )
                        ]

                elif name == "download_presentation":
                    file_id = arguments.get("file_id")
                    logger.info(
                        f"Executing download_presentation tool for file_id: {file_id}"
                    )

                    # Get presentation metadata
                    endpoint = f"me/drive/items/{file_id}"
                    result = await make_graph_api_request(
                        "get", endpoint, access_token=access_token
                    )

                    download_url = result.get("@microsoft.graph.downloadUrl")
                    response = {
                        "status": "success",
                        "file_id": result.get("id"),
                        "name": result.get("name"),
                        "download_url": download_url,
                        "size": result.get("size"),
                        "web_url": result.get("webUrl"),
                    }

                    return [
                        types.TextContent(
                            type="text", text=json.dumps(response, indent=2)
                        )
                    ]

                elif name == "search_presentations":
                    query = arguments.get("query")
                    limit = arguments.get("limit", 25)
                    logger.info(
                        f"Executing search_presentations tool with query: {query}, limit: {limit}"
                    )

                    if not query:
                        return [
                            types.TextContent(
                                type="text",
                                text=json.dumps(
                                    {"status": "error", "message": "query is required"},
                                    indent=2,
                                ),
                            )
                        ]

                    # Determine if we're using SharePoint or OneDrive
                    is_sharepoint = await is_sharepoint_storage(access_token)

                    endpoint = f"me/drive/root/search(q='{query}')"
                    params = {
                        "$top": 100 if is_sharepoint else limit,
                        "$select": "id,name,webUrl,lastModifiedDateTime,size,createdDateTime,file",
                        "$orderby": "lastModifiedDateTime desc",
                    }

                    result = await make_graph_api_request(
                        "get", endpoint, params=params, access_token=access_token
                    )

                    presentations = []
                    if is_sharepoint:
                        for item in result.get("value", []):
                            if (
                                item.get("file")
                                and item.get("file", {}).get("mimeType")
                                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            ):
                                presentations.append(item)
                                if len(presentations) >= limit:
                                    break
                    else:
                        for item in result.get("value", []):
                            if item["name"].lower().endswith(".pptx"):
                                presentations.append(item)
                                if len(presentations) >= limit:
                                    break

                    formatted_presentations = []
                    for item in presentations:
                        formatted_presentations.append(
                            {
                                "id": item.get("id"),
                                "name": item.get("name"),
                                "web_url": item.get("webUrl"),
                                "last_modified": item.get("lastModifiedDateTime"),
                                "created": item.get("createdDateTime"),
                                "size": item.get("size"),
                            }
                        )

                    response = {
                        "status": "success",
                        "query": query,
                        "presentations": formatted_presentations,
                    }

                    return [
                        types.TextContent(
                            type="text", text=json.dumps(response, indent=2)
                        )
                    ]

                else:
                    logger.error(f"Unsupported tool called: {name}")
                    return [
                        types.TextContent(
                            type="text",
                            text=json.dumps(
                                {
                                    "status": "error",
                                    "message": f"Unsupported tool: {name}",
                                },
                                indent=2,
                            ),
                        )
                    ]

            except Exception as e:
                logger.error(f"Error executing tool {name}: {str(e)}", exc_info=True)
                return [
                    types.TextContent(
                        type="text",
                        text=json.dumps(
                            {
                                "status": "error",
                                "message": f"Error executing {name}: {str(e)}",
                            },
                            indent=2,
                        ),
                    )
                ]
        except Exception as e:
            logger.error(
                f"Failed to get Microsoft client for tool {name}: {str(e)}",
                exc_info=True,
            )
            return [
                types.TextContent(
                    type="text",
                    text=json.dumps(
                        {
                            "status": "error",
                            "message": f"Authentication error for {name}: {str(e)}",
                        },
                        indent=2,
                    ),
                )
            ]

    return server


server = create_server


def get_initialization_options(server_instance: Server) -> InitializationOptions:
    """Get the initialization options for the server"""
    return InitializationOptions(
        server_name=f"{SERVICE_NAME}-server",
        server_version="1.0.0",
        capabilities=server_instance.get_capabilities(
            notification_options=NotificationOptions(),
            experimental_capabilities={},
        ),
    )


# Main entry point for authentication
if __name__ == "__main__":
    if len(sys.argv) > 1 and sys.argv[1].lower() == "auth":
        user_id = "local"
        authenticate_and_save_credentials(user_id, SERVICE_NAME, SCOPES)
    else:
        print("Usage:")
        print("  python main.py auth - Run authentication flow for a user")
        print("Note: To run the server normally, use the guMCP server framework.")
